"""
Build an extended version of the agentic.pptx presentation.
Adds ~25 new slides covering LangChain workflows, error handling patterns,
and the hands-on lab exercise — bridging the theory in the original deck
to practical implementation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Open existing presentation ──
prs = Presentation("agentic.pptx")

# ── Helper: find layouts by name ──
layouts = {layout.name: layout for layout in prs.slide_master.slide_layouts}
print("Available layouts:", list(layouts.keys()))

# Pick common layouts
title_layout = layouts.get("Title Slide") or layouts.get("White Title Slide")
content_layout = layouts.get("Title and Content")
blank_layout = layouts.get("Blank")

# ── Style constants (match existing deck) ──
DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
RED = RGBColor(0xFF, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)
ACCENT_GREEN = RGBColor(0x00, 0x80, 0x00)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)


def add_title_slide(title_text, subtitle_text=""):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(title_layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:  # Title
            shape.text = title_text
        elif shape.placeholder_format.idx == 1 and subtitle_text:  # Subtitle
            shape.text = subtitle_text
    return slide


def add_content_slide(title_text, bullets, sub_bullets=None):
    """Add a title + bullet points slide. sub_bullets is a dict mapping bullet index to list of sub-items."""
    slide = prs.slides.add_slide(content_layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:  # Title
            shape.text = title_text
        elif shape.placeholder_format.idx == 1:  # Content
            tf = shape.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = Pt(6)
                for run in p.runs:
                    run.font.size = Pt(20)

                # Add sub-bullets if present
                if sub_bullets and i in sub_bullets:
                    for sub in sub_bullets[i]:
                        sp = tf.add_paragraph()
                        sp.text = sub
                        sp.level = 1
                        sp.space_after = Pt(4)
                        for run in sp.runs:
                            run.font.size = Pt(16)
                            run.font.color.rgb = MEDIUM_GRAY
    return slide


def add_code_slide(title_text, code_text, language_label="", note_text=""):
    """Add a slide with a code block (dark background text box)."""
    slide = prs.slides.add_slide(content_layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title_text
        elif shape.placeholder_format.idx == 1:
            # Remove default content placeholder text
            shape.text_frame.clear()

    # Code box
    left = Inches(0.8)
    top = Inches(1.8) if not note_text else Inches(1.6)
    width = Inches(10.5)
    code_height = Inches(4.0) if not note_text else Inches(3.5)

    # Dark background rectangle
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, code_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    rect.line.fill.background()

    # Code text
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    if language_label:
        p = tf.paragraphs[0]
        p.text = language_label
        run = p.runs[0]
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        run.font.bold = True
        p.space_after = Pt(4)
        # Add code after label
        for line in code_text.split("\n"):
            cp = tf.add_paragraph()
            cp.text = line
            for r in cp.runs:
                r.font.size = Pt(13)
                r.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
                r.font.name = "Courier New"
            cp.space_after = Pt(1)
    else:
        tf.paragraphs[0].text = ""
        for j, line in enumerate(code_text.split("\n")):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
                r.font.name = "Courier New"
            p.space_after = Pt(1)

    # Optional note below code
    if note_text:
        note_top = top + code_height + Inches(0.15)
        note_box = slide.shapes.add_textbox(left, note_top, width, Inches(0.8))
        ntf = note_box.text_frame
        ntf.word_wrap = True
        np = ntf.paragraphs[0]
        np.text = note_text
        for r in np.runs:
            r.font.size = Pt(14)
            r.font.color.rgb = MEDIUM_GRAY
            r.font.italic = True

    return slide


def add_two_column_slide(title_text, left_title, left_bullets, right_title, right_bullets):
    """Add a slide with two columns of content."""
    slide = prs.slides.add_slide(blank_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10.5), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK_GRAY

    # Left column
    col_top = Inches(1.5)
    col_width = Inches(5.0)

    # Left title
    ltb = slide.shapes.add_textbox(Inches(0.8), col_top, col_width, Inches(0.5))
    lp = ltb.text_frame.paragraphs[0]
    lp.text = left_title
    lp.runs[0].font.size = Pt(22)
    lp.runs[0].font.bold = True
    lp.runs[0].font.color.rgb = ACCENT_BLUE

    # Left bullets
    lbb = slide.shapes.add_textbox(Inches(0.8), col_top + Inches(0.6), col_width, Inches(4.5))
    ltf = lbb.text_frame
    ltf.word_wrap = True
    for i, b in enumerate(left_bullets):
        if i == 0:
            p = ltf.paragraphs[0]
        else:
            p = ltf.add_paragraph()
        p.text = f"\u2022 {b}"
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = Pt(16)

    # Right column
    rtb = slide.shapes.add_textbox(Inches(6.3), col_top, col_width, Inches(0.5))
    rp = rtb.text_frame.paragraphs[0]
    rp.text = right_title
    rp.runs[0].font.size = Pt(22)
    rp.runs[0].font.bold = True
    rp.runs[0].font.color.rgb = ACCENT_BLUE

    rbb = slide.shapes.add_textbox(Inches(6.3), col_top + Inches(0.6), col_width, Inches(4.5))
    rtf = rbb.text_frame
    rtf.word_wrap = True
    for i, b in enumerate(right_bullets):
        if i == 0:
            p = rtf.paragraphs[0]
        else:
            p = rtf.add_paragraph()
        p.text = f"\u2022 {b}"
        p.space_after = Pt(6)
        for r in p.runs:
            r.font.size = Pt(16)

    return slide


def add_diagram_slide(title_text, boxes, note_text=""):
    """Add a simple flow diagram with connected boxes.
    boxes: list of (label, color) tuples, drawn left-to-right with arrows.
    """
    slide = prs.slides.add_slide(blank_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10.5), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK_GRAY

    n = len(boxes)
    box_width = Inches(1.8)
    box_height = Inches(1.0)
    arrow_width = Inches(0.4)
    total_width = n * box_width + (n - 1) * arrow_width
    start_x = (Inches(12.19) - total_width) / 2
    center_y = Inches(3.2)

    for i, (label, color) in enumerate(boxes):
        x = start_x + i * (box_width + arrow_width)
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(center_y), int(box_width), int(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        for r in p.runs:
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = WHITE

        # Arrow (except after last box)
        if i < n - 1:
            arrow_x = x + box_width
            arrow_y = center_y + box_height / 2 - Inches(0.15)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, int(arrow_x), int(arrow_y), int(arrow_width), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MEDIUM_GRAY
            arrow.line.fill.background()

    # Note
    if note_text:
        nb = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(10.5), Inches(1.2))
        ntf = nb.text_frame
        ntf.word_wrap = True
        np = ntf.paragraphs[0]
        np.text = note_text
        np.alignment = PP_ALIGN.CENTER
        for r in np.runs:
            r.font.size = Pt(16)
            r.font.color.rgb = MEDIUM_GRAY
            r.font.italic = True

    return slide


def add_table_slide(title_text, headers, rows):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(blank_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10.5), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK_GRAY

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(11.2)
    height = Inches(0.45) * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    col_width = int(width / n_cols)
    for col in table.columns:
        col.width = col_width

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7)

    return slide


# ════════════════════════════════════════════════════════════
# NEW SLIDES — Section: Building AI Workflows with LangChain
# ════════════════════════════════════════════════════════════

# ── Section Divider ──
add_title_slide(
    "Building AI Workflows with LangChain",
    "From Theory to Practice: Hands-On Lab"
)

# ── What is LangChain? ──
add_content_slide(
    "What is LangChain?",
    [
        "A framework for building applications powered by LLMs",
        "Available in Python (langchain) and JavaScript (langchain.js)",
        "Provides composable building blocks:",
        "Integrates with 50+ LLM providers (OpenAI, Anthropic, etc.)",
    ],
    sub_bullets={
        2: [
            "Prompt Templates — reusable, parameterized prompts",
            "Chains — connect components with the pipe (|) operator",
            "Output Parsers — structured data from LLM text",
            "Tools — let LLMs call functions and APIs",
        ]
    }
)

# ── What is a Workflow/Chain? ──
add_diagram_slide(
    "What is a Workflow (Chain)?",
    [
        ("Input\n(Topic)", ACCENT_BLUE),
        ("Step 1\nGenerate\nOutline", ACCENT_GREEN),
        ("Step 2\nWrite\nDraft", ACCENT_ORANGE),
        ("Step 3\nSummarize", RGBColor(0x80, 0x00, 0x80)),
        ("Output\n(Summary)", ACCENT_BLUE),
    ],
    "Each step can be an LLM call, data transformation, API call, or validation check."
)

# ── Why Workflows? ──
add_content_slide(
    "Why Use Multi-Step Workflows?",
    [
        "Decomposition — break complex tasks into manageable steps",
        "Quality — each step can be specialized and optimized",
        "Debuggability — inspect output at each stage",
        "Reusability — individual steps work across different pipelines",
        "Error handling — failures at one step don't crash everything",
    ]
)

# ── Workflow Patterns ──
add_two_column_slide(
    "Common Workflow Patterns",
    "Sequential Chain",
    [
        "Output of step N feeds step N+1",
        "Simplest pattern, easy to debug",
        "Example: Topic → Outline → Draft → Summary",
        "Use when steps have clear dependencies",
    ],
    "Parallel Execution",
    [
        "Run independent steps concurrently",
        "Combine results when all complete",
        "Faster: N parallel calls ≈ 1 call latency",
        "Example: Answer 3 sub-questions at once",
    ]
)

add_two_column_slide(
    "More Workflow Patterns",
    "Conditional Branching",
    [
        "LLM classifies input → different paths",
        "Simple questions get direct answers",
        "Complex questions get decomposed",
        "Saves cost on easy inputs",
    ],
    "Fallback Chains",
    [
        "If primary model/approach fails, try backup",
        "GPT-4 fails? Try GPT-3.5 as fallback",
        "Structured output fails? Try simpler format",
        "Graceful degradation > hard failure",
    ]
)

# ── LangChain Pipe Syntax ──
add_code_slide(
    "LangChain: The Pipe Operator",
    """# Python — compose with the | operator
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
import os

# xAI Grok — free API, OpenAI-compatible
llm = ChatOpenAI(
    model="grok-3-mini-fast",
    base_url="https://api.x.ai/v1",
    api_key=os.environ["XAI_API_KEY"],
)

chain = (
    ChatPromptTemplate.from_template("Tell me about {topic}")
    | llm
    | StrOutputParser()
)
result = chain.invoke({"topic": "WebSockets"})""",
    "PYTHON",
    "The pipe operator connects Prompt → LLM → Output Parser into a single callable chain."
)

add_code_slide(
    "LangChain.js: Same Pattern in JavaScript",
    """// JavaScript — same pipe pattern
import { ChatOpenAI } from "@langchain/openai";
import { ChatPromptTemplate } from "@langchain/core/prompts";
import { StringOutputParser } from "@langchain/core/output_parsers";

// xAI Grok — free API, OpenAI-compatible
const llm = new ChatOpenAI({
    model: "grok-3-mini-fast",
    configuration: {
        baseURL: "https://api.x.ai/v1",
        apiKey: process.env.XAI_API_KEY,
    },
});

const chain = ChatPromptTemplate
    .fromTemplate("Tell me about {topic}")
    .pipe(llm)
    .pipe(new StringOutputParser());

const result = await chain.invoke({ topic: "WebSockets" });""",
    "JAVASCRIPT",
    "LangChain.js uses .pipe() method — same concept, JavaScript syntax."
)

# ── Multi-Step Workflow Code ──
add_code_slide(
    "Building a Multi-Step Workflow",
    """# Define 3 separate chains
outline_chain = outline_prompt | llm | StrOutputParser()
draft_chain   = draft_prompt   | llm | StrOutputParser()
summary_chain = summary_prompt | llm | StrOutputParser()

# Run step by step (see intermediate results)
outline = outline_chain.invoke({"topic": topic})
draft   = draft_chain.invoke({"outline": outline})
summary = summary_chain.invoke({"draft": draft})

# OR compose into a single chain
full_workflow = (
    {"outline": outline_chain}
    | draft_chain
    | RunnableLambda(lambda draft: {"draft": draft})
    | summary_chain
)""",
    "PYTHON",
    "Step-by-step execution lets you inspect intermediates. Composed chains run as one call."
)

# ══════════════════════════════════════
# Error Handling Section
# ══════════════════════════════════════

add_title_slide(
    "Error Handling in AI Workflows",
    "The Most Important Part You'll Want to Skip"
)

# ── Error Types Table ──
add_table_slide(
    "Types of Errors in LLM Workflows",
    ["Error Type", "What Happens", "Example", "Strategy"],
    [
        ["API Errors", "Rate limits, timeouts, auth failures", "429 Too Many Requests", "Retry with exponential backoff"],
        ["Malformed Output", "LLM returns unexpected format", "Asked for JSON, got prose", "Output parsers + validation"],
        ["Token Limits", "Input/output exceeds context window", "Prompt too long for model", "Truncate, chunk, or use larger model"],
        ["Hallucination", "LLM generates incorrect content", "Cites non-existent sources", "Validation steps, guardrails"],
        ["Chain Propagation", "Error in step N corrupts step N+1", "Bad outline → bad draft", "Validate intermediate outputs"],
        ["Cost Explosion", "Infinite loops or excessive retries", "Agent calls API 1000 times", "Step limits, budget caps"],
    ]
)

# ── Retry Pattern ──
add_code_slide(
    "Pattern 1: Retry with Exponential Backoff",
    """async function callWithRetry(fn, maxRetries = 3) {
    for (let attempt = 1; attempt <= maxRetries; attempt++) {
        try {
            return await fn();
        } catch (error) {
            const isRetryable = error.status === 429
                             || error.status >= 500;
            if (isRetryable && attempt < maxRetries) {
                const delay = Math.pow(2, attempt) * 1000;
                console.log(`Retry ${attempt} in ${delay}ms...`);
                await new Promise(r => setTimeout(r, delay));
            } else {
                throw error;  // Non-retryable or max retries hit
            }
        }
    }
}""",
    "JAVASCRIPT",
    "Rate limits (429) and server errors (5xx) are retryable. Auth errors (401) are not."
)

# ── Validation Pattern ──
add_code_slide(
    "Pattern 2: Output Validation",
    """def validate_json_output(llm_output, required_keys):
    \"\"\"Validate LLM output is valid JSON with required keys.\"\"\"
    text = llm_output.strip()

    # Extract from markdown code fences if present
    if "```json" in text:
        start = text.index("```json") + 7
        end = text.index("```", start)
        text = text[start:end].strip()

    try:
        parsed = json.loads(text)
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON: {e}")

    missing = [k for k in required_keys if k not in parsed]
    if missing:
        raise ValueError(f"Missing keys: {missing}")

    return parsed""",
    "PYTHON",
    "Never trust raw LLM output. Always parse and validate before passing to the next step."
)

# ── Fallback Pattern ──
add_code_slide(
    "Pattern 3: Fallback Chains",
    """# LangChain built-in fallbacks
primary_llm  = ChatOpenAI(model="grok-3-mini-fast", ...)
fallback_llm = ChatOpenAI(model="grok-3-mini-fast", temperature=0.3, ...)

# If primary fails, automatically try fallback
robust_llm = primary_llm.with_fallbacks([fallback_llm])

# Custom fallback wrapper
async function withFallback(primaryFn, fallbackFn) {
    try {
        return await primaryFn();
    } catch (error) {
        console.warn("Primary failed, using fallback");
        return await fallbackFn();
    }
}""",
    "PYTHON + JS",
    "Fallbacks provide graceful degradation: a slightly worse answer beats a crash."
)

# ── Mapping Errors to the 5 Stages ──
add_content_slide(
    "Error Handling Across the 5 Stages",
    [
        "Stage 1 (Simple Flows): Retry + validate output format",
        "Stage 2 (RAG): Validate retrieval quality, handle empty results",
        "Stage 3 (Agentic): Step limits, cost caps, loop detection",
        "Stage 4 (Production): Circuit breakers, observability, alerting",
        "Stage 5 (Enterprise): RBAC, audit trails, compliance checks",
    ],
    sub_bullets={
        0: ["Your code should handle this today"],
        1: ["Check that retrieved context is relevant before sending to LLM"],
        2: ["The most dangerous stage — autonomy introduces side effects"],
        3: ["Treat AI calls like any external service: monitor, log, alert"],
        4: ["Security and governance become the primary failure surface"],
    }
)

# ── Real-World Failure Example ──
add_content_slide(
    'Real-World Failure: "The JSON That Wasn\'t"',
    [
        "Prompt: \"Return a JSON object with title, tags, and difficulty\"",
        "LLM returns: \"Sure! Here's the JSON:\\n```json\\n{...}\\n```\"",
        "Your code calls JSON.parse() on the full response...",
        "Result: SyntaxError crashes the entire workflow",
    ],
    sub_bullets={
        3: [
            "Fix: Strip markdown fences, extract JSON block, then parse",
            "Better fix: Use LangChain's StructuredOutputParser",
            "Best fix: Use the model's structured output / JSON mode",
        ]
    }
)

# ══════════════════════════════════════
# Lab Exercise Section
# ══════════════════════════════════════

add_title_slide(
    "Hands-On Lab: LangChain Workflows",
    "Build, Break, and Fix AI Pipelines"
)

# ── Lab Overview ──
add_content_slide(
    "Lab Structure",
    [
        "Part 1: Python Workflow (Jupyter Notebook) — ~30 min",
        "Part 2: Node.js Workflows (JavaScript) — ~30 min",
        "Part 3: Your Novel Workflow — ~30 min",
    ],
    sub_bullets={
        0: [
            "Simple chain → Multi-step pipeline → Error handling",
            "Retry logic, output validation, fallback chains",
        ],
        1: [
            "node-workflow-basic/: Sequential content generation pipeline",
            "node-workflow-advanced/: Parallel execution, branching, structured output",
        ],
        2: [
            "Design your own 3+ step workflow",
            "Must include error handling",
            "Create a 1-slide summary",
        ],
    }
)

# ── Lab Directory Structure ──
add_code_slide(
    "Lab Repository Structure",
    """Agents1/
├── README.md                           # Full lab instructions
├── .devcontainer/devcontainer.json     # Codespaces auto-setup
├── python-workflow/
│   ├── requirements.txt
│   └── workflow_example.ipynb          # Jupyter notebook walkthrough
├── node-workflow-basic/
│   ├── package.json
│   └── workflow.js                     # Sequential pipeline
├── node-workflow-advanced/
│   ├── package.json
│   └── workflow.js                     # Parallel + branching + fallbacks
└── my-workflow/                        # YOUR NOVEL WORKFLOW
    ├── workflow.[py|js|ipynb]
    └── slide.md                        # 1-slide summary""",
    "REPOSITORY LAYOUT",
)

# ── Python Example Walkthrough ──
add_content_slide(
    "Python Notebook Walkthrough",
    [
        "Cell 1-2: Setup — install deps, configure API key",
        "Cell 3: Simple chain — prompt | llm | output_parser",
        "Cell 4: Multi-step — outline → draft → summary (see intermediates)",
        "Cell 5: Composed chain — single .invoke() runs all 3 steps",
        "Cell 6: Retry with exponential backoff",
        "Cell 7: JSON output validation (extract from code fences)",
        "Cell 8: LangChain .with_fallbacks() for graceful degradation",
        "Cell 9: Full robust workflow combining all patterns",
    ]
)

# ── Node Basic Example ──
add_diagram_slide(
    "Node Basic: Sequential Content Pipeline",
    [
        ("Topic\n(Input)", ACCENT_BLUE),
        ("Generate\nOutline", ACCENT_GREEN),
        ("Validate\nLength", ACCENT_ORANGE),
        ("Write\nDraft", ACCENT_GREEN),
        ("Create\nSummary", ACCENT_GREEN),
    ],
    "Each step includes try/catch with retry. Validation between steps catches bad intermediate output."
)

# ── Node Advanced Example ──
add_content_slide(
    "Node Advanced: Research Assistant",
    [
        "Step 1: Classify question as simple or complex",
        "Simple path: Direct answer (fast, cheap)",
        "Complex path: Decompose → Parallel answer → Synthesize",
        "Step 4: Format output as validated JSON",
        "Error handling at every level:",
    ],
    sub_bullets={
        4: [
            "Custom error classes (RetryableError, ValidationError)",
            "Fallback sub-question generation if JSON parsing fails",
            "Promise.all for parallel execution with per-item error catching",
            "Unstructured fallback if final formatting fails",
        ]
    }
)

# ── Advanced: Conditional Branching Diagram ──
add_code_slide(
    "Advanced: Conditional Branching Logic",
    """// Step 1: Classify the question
const classification = await classifyQuestion(question);

if (classification === "simple") {
    // Fast path — single LLM call
    answer = await directAnswerChain.invoke({ question });
} else {
    // Complex path — decompose, parallel answer, synthesize
    const subQuestions = await generateSubQuestions(question);

    // Answer ALL sub-questions in parallel
    const answers = await Promise.all(
        subQuestions.map(sq => answerChain.invoke({ sq }))
    );

    answer = await synthesizeChain.invoke({
        originalQuestion: question,
        subAnswers: answers
    });
}""",
    "JAVASCRIPT",
    "Branching saves cost and latency on simple inputs. Parallel execution handles complex ones efficiently."
)

# ══════════════════════════════════════
# Student Project Section
# ══════════════════════════════════════

add_title_slide(
    "Your Novel Workflow",
    "Design Something Original"
)

# ── Requirements ──
add_content_slide(
    "Student Project Requirements",
    [
        "At least 3 steps (LLM calls, transformations, or API calls)",
        "Error handling (try/catch + validate at least one intermediate output)",
        "Something novel — not a copy of the provided examples",
        "1-slide summary showing workflow, application context, and error strategy",
    ],
    sub_bullets={
        2: [
            "Think: what multi-step AI task would be useful in a real web app?",
        ]
    }
)

# ── Project Ideas ──
add_two_column_slide(
    "Project Ideas",
    "Content & Text",
    [
        "Code Review Pipeline: code → identify issues → suggest fixes → improved version",
        "Study Guide Creator: notes → key concepts → practice questions → summary",
        "Email Composer: key points + tone → draft → check issues → final",
        "Meeting Summarizer: raw notes → action items → prioritize → follow-up email",
    ],
    "Data & Web",
    [
        "Recipe Generator: ingredients → suggest recipes → full recipe + nutrition",
        "Product Description Writer: specs → descriptions for audiences → A/B variants",
        "Accessibility Checker: HTML → issues → fixes → improved HTML",
        "Data Pipeline Describer: SQL schema → docs → sample queries → explain",
    ]
)

# ── Slide Template ──
add_content_slide(
    "Your 1-Slide Summary Should Include",
    [
        "What does your workflow do? (1-2 sentences)",
        "What are the steps? (diagram or numbered list)",
        "What web application would use this? (broader context)",
        "What errors can occur and how do you handle them?",
    ],
    sub_bullets={
        1: [
            "Show the data flow: Input → Step 1 → Step 2 → ... → Output"
        ],
        2: [
            "E.g., \"Part of an e-commerce platform that auto-generates product listings\""
        ],
        3: [
            "E.g., \"Retry on rate limits, validate JSON output, fallback to simpler model\""
        ],
    }
)

# ══════════════════════════════════════
# Connecting Theory to Practice
# ══════════════════════════════════════

add_title_slide(
    "Connecting the Dots",
    "From Enterprise AI Systems to Your Lab Code"
)

add_content_slide(
    "Your Lab Code Maps to Enterprise Stages",
    [
        "Your basic chain = Stage 1 (Simple LLM Flows)",
        "Your sequential pipeline = Stage 2 (RAG-style multi-step)",
        "Your branching workflow = Stage 3 (Agentic decisions)",
        "Adding validation + retry + fallbacks = Stage 4 (Production-ready)",
        "RBAC + audit + multi-tenant = Stage 5 (where your career goes next)",
    ],
    sub_bullets={
        0: ["Prompt → LLM → Response. Failure mode: hallucination."],
        1: ["Multiple steps, each with potential failure points."],
        2: ["LLM decides which path to take. Failure mode: wrong action."],
        3: ["Error handling makes your code deployable."],
        4: ["The patterns you learn today scale to enterprise systems."],
    }
)

# ── Security Considerations ──
add_content_slide(
    "Security: Prompt Injection in Workflows",
    [
        "What if malicious input in Step 1 affects Step 3?",
        "Example: user input includes \"Ignore previous instructions and...\"",
        "In a chain, this gets passed through as context to later steps",
        "Mitigations:",
    ],
    sub_bullets={
        3: [
            "Sanitize user input before feeding to chains",
            "Use separate system prompts that can't be overridden",
            "Validate intermediate outputs (does the outline look like an outline?)",
            "Never let LLM output be executed as code without review",
        ]
    }
)

# ── Testing AI Workflows ──
add_content_slide(
    "Testing AI Workflows",
    [
        "Challenge: LLMs are non-deterministic — same input, different output",
        "Strategy 1: Set temperature=0 for reproducible tests",
        "Strategy 2: Test structure, not exact content",
        "Strategy 3: Use LLM-as-judge for quality assertions",
        "Strategy 4: Mock LLM responses for unit tests, use real LLM for integration",
    ],
    sub_bullets={
        2: [
            "Does the JSON have the right keys? Is the outline > 50 chars? Is confidence one of 3 values?"
        ],
        3: [
            "\"Does this summary accurately reflect the draft?\" → LLM grades the output"
        ],
    }
)

# ── Cost Awareness ──
add_content_slide(
    "Cost Awareness in AI Workflows",
    [
        "Every LLM call costs money (tokens in + tokens out)",
        "A 3-step workflow = 3x the cost of a single call",
        "Parallel sub-questions can multiply costs quickly",
        "Retries add up — 3 retries × 3 steps = up to 9 calls",
        "Best practices:",
    ],
    sub_bullets={
        4: [
            "Use faster/cheaper models (grok-3-mini-fast) for classification/simple steps",
            "Cache results where possible",
            "Set max_retries and budget limits",
            "Log token usage to understand costs",
        ]
    }
)

# ══════════════════════════════════════
# Wrap-Up
# ══════════════════════════════════════

add_title_slide(
    "Key Takeaways",
    ""
)

add_content_slide(
    "What You Should Take Away",
    [
        "AI workflows decompose complex tasks into manageable, debuggable steps",
        "LangChain provides composable building blocks (prompts, chains, parsers)",
        "Error handling is not optional — LLMs fail in unique ways",
        "The patterns you learn today (retry, validate, fallback) scale to production",
        "Competitive advantage comes from system design, not just prompts",
    ]
)

add_content_slide(
    "Discussion Questions",
    [
        "When should you use a multi-step workflow vs. a single prompt?",
        "How do you decide what to retry vs. what to fail?",
        "What are the security implications of chaining LLM outputs?",
        "How would you test a non-deterministic AI workflow?",
        "Where does error handling belong — each step, top level, or both?",
    ]
)

add_content_slide(
    "Next Steps",
    [
        "Complete the lab — work through Python notebook and Node.js examples",
        "Build your novel workflow (3+ steps, error handling, 1-slide summary)",
        "Submit: push your my-workflow/ directory to your assignment repo",
        "Come to class ready to present your 1-slide summary",
    ]
)

# ── Save ──
output_path = "agentic_extended.pptx"
prs.save(output_path)
print(f"\nSaved extended presentation to {output_path}")
print(f"Total slides: {len(prs.slides)}")
